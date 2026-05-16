import os
from pypdf import PdfReader
from uuid import UUID
from langchain_text_splitters import RecursiveCharacterTextSplitter
from qdrant_client.models import PointStruct, VectorParams, Distance
from app.core.llm import get_qdrant_client, get_embedding_model
from app.config import settings

COLLECTION_NAME = "study_documents"

def _ensure_collection_exists():
    client = get_qdrant_client()
    try:
        client.get_collection(COLLECTION_NAME)
    except Exception:
        # Create collection if not exists (embedding dim 384 for all-MiniLM-L6-v2)
        client.create_collection(
            collection_name=COLLECTION_NAME,
            vectors_config=VectorParams(size=384, distance=Distance.COSINE)
        )

def process_and_index_document(file_path: str, document_id: UUID, user_id: UUID, filename: str) -> int:
    """Extracts text, chunks it, embeds it, and stores in Qdrant. Returns chunk_count."""
    _ensure_collection_exists()
    
    # 1. Extract Text
    text = ""
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        elif ext == ".docx":
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
        else:
            # Fallback: try plain text
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
    except Exception as e:
        print(f"Error reading file {file_path} (ext={ext}): {e}")
        try:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        except Exception as e2:
            raise Exception(f"Could not extract text from document: {e2}")

    if not text.strip():
        raise Exception("Document is empty or text could not be extracted")

    # 2. Chunking
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    if not chunks:
        return 0

    # 3. Embedding
    model = get_embedding_model()
    embeddings = model.encode(chunks, show_progress_bar=False)

    # 4. Store in Qdrant
    client = get_qdrant_client()
    points = []
    import uuid
    for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        points.append(
            PointStruct(
                id=str(uuid.uuid4()),
                vector=embedding.tolist(),
                payload={
                    "document_id": str(document_id),
                    "user_id": str(user_id),
                    "filename": filename,
                    "text": chunk,
                    "chunk_index": i
                }
            )
        )
    
    # Upload points in batches
    client.upload_points(
        collection_name=COLLECTION_NAME,
        points=points
    )
    
    return len(chunks)

def delete_document_vectors(document_id: UUID):
    """Deletes all vectors associated with a document_id."""
    try:
        client = get_qdrant_client()
        from qdrant_client.http.models import Filter, FieldCondition, MatchValue
        client.delete(
            collection_name=COLLECTION_NAME,
            points_selector=Filter(
                must=[
                    FieldCondition(
                        key="document_id",
                        match=MatchValue(value=str(document_id))
                    )
                ]
            )
        )
    except Exception:
        pass

def search_context(user_id: UUID, query: str, top_k: int = 3) -> list[str]:
    """Searches for relevant chunks for a specific user."""
    _ensure_collection_exists()
    try:
        model = get_embedding_model()
        query_vector = model.encode(query).tolist()
        
        client = get_qdrant_client()
        from qdrant_client.http.models import Filter, FieldCondition, MatchValue
        
        results = client.query_points(
            collection_name=COLLECTION_NAME,
            query=query_vector,
            query_filter=Filter(
                must=[
                    FieldCondition(
                        key="user_id",
                        match=MatchValue(value=str(user_id))
                    )
                ]
            ),
            limit=top_k
        )
        
        return [res.payload["text"] for res in results.points if res.payload]
    except Exception as e:
        print(f"Error searching context: {e}")
        return []
