import os
from uuid import UUID
from sqlalchemy.orm import Session
from app.models.document import Document
from app.services.rag_service import process_and_index_document
from app.core.llm import get_llm
from fastapi import HTTPException

def process_document_background(db: Session, document_id: UUID, file_path: str, user_id: UUID, filename: str):
    """Background task to extract, embed and store document chunks."""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        return
        
    try:
        chunk_count = process_and_index_document(file_path, document_id, user_id, filename)
        document.status = "ready"
        document.chunk_count = chunk_count
    except Exception as e:
        print(f"Error processing document {document_id}: {e}")
        document.status = "failed"
        
    db.commit()

def generate_document_summary(db: Session, document: Document) -> str:
    """Summarizes document by extracting text using PyMuPDF and passing to LLM."""
    if not document.file_path or not os.path.exists(document.file_path):
        raise HTTPException(status_code=404, detail="File not found on server")
        
    text = ""
    ext = os.path.splitext(document.file_path)[1].lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(document.file_path)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        elif ext == ".docx":
            from docx import Document as DocxDocument
            doc = DocxDocument(document.file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(document.file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
        else:
            with open(document.file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
    except Exception:
        try:
            with open(document.file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        except:
            raise HTTPException(status_code=500, detail="Failed to read document content for summarization")
            
    # Truncate text if too long (Gemini context is large, but to be safe and fast)
    text = text[:30000] 
    
    prompt = f"""Hãy tóm tắt tài liệu sau đây trong 3-5 đoạn văn bằng TIẾNG VIỆT.
Tập trung vào các khái niệm chính, ý chính và các chi tiết quan trọng.
Hãy viết rõ ràng, súc tích, phù hợp cho mục đích học tập.

Nội dung tài liệu:
{text}"""

    llm = get_llm(temperature=0.3)
    response = llm.invoke([("human", prompt)])
    
    content_raw = response.content
    if isinstance(content_raw, list):
        summary = "".join([c.get("text", "") for c in content_raw if isinstance(c, dict) and "text" in c])
    else:
        summary = str(content_raw)
    summary = summary.strip()
    
    document.summary = summary
    db.commit()
    
    return summary
