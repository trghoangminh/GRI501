import os
import shutil
import mimetypes
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Query, BackgroundTasks
from sqlalchemy.orm import Session
from uuid import UUID
from typing import List, Optional
from app.database import get_db
from app.models.user import User
from app.models.document import Document
from app.schemas.document import DocumentResponse, DocumentSummaryResponse
from app.core.dependencies import get_current_user
from app.config import settings
from app.services.document_service import process_document_background, generate_document_summary
from app.services.quiz_service import generate_quiz
from app.services.rag_service import delete_document_vectors

router = APIRouter(prefix="/documents", tags=["documents"])

ALLOWED_MIME_TYPES = [
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document", # DOCX
    "text/plain",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation" # PPTX
]

@router.post("/upload", response_model=DocumentResponse)
async def upload_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...), 
    db: Session = Depends(get_db), 
    current_user: User = Depends(get_current_user)
):
    # Check size
    file.file.seek(0, 2)
    file_size = file.file.tell()
    file.file.seek(0)
    
    if file_size > settings.MAX_FILE_SIZE_MB * 1024 * 1024:
        raise HTTPException(status_code=413, detail=f"File too large. Max {settings.MAX_FILE_SIZE_MB}MB")
        
    # Check mime type
    mime = mimetypes.guess_type(file.filename)[0] or "application/octet-stream"
    
    if mime not in ALLOWED_MIME_TYPES:
        raise HTTPException(status_code=415, detail=f"Unsupported file type: {mime}")
        
    ext = file.filename.split(".")[-1] if "." in file.filename else "bin"
    file_type = ext.lower()
    
    docs_dir = os.path.join(settings.UPLOAD_DIR, "documents")
    os.makedirs(docs_dir, exist_ok=True)
    
    import uuid
    new_filename = f"{current_user.id}_{uuid.uuid4().hex[:8]}.{ext}"
    file_path = os.path.join(docs_dir, new_filename)
    
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)
        
    document = Document(
        user_id=current_user.id,
        filename=new_filename,
        original_name=file.filename,
        file_type=file_type,
        file_path=file_path,
        file_size=file_size,
        status="processing"
    )
    db.add(document)
    db.commit()
    db.refresh(document)
    
    # Background processing for RAG
    background_tasks.add_task(process_document_background, db, document.id, file_path, current_user.id, file.filename)
    
    return document

@router.get("/", response_model=List[DocumentResponse])
def get_documents(
    type: Optional[str] = None,
    status: Optional[str] = None,
    search: Optional[str] = None,
    db: Session = Depends(get_db), 
    current_user: User = Depends(get_current_user)
):
    query = db.query(Document).filter(Document.user_id == current_user.id)
    if type:
        query = query.filter(Document.file_type == type.lower())
    if status:
        query = query.filter(Document.status == status)
    if search:
        query = query.filter(Document.original_name.ilike(f"%{search}%"))
        
    return query.order_by(Document.created_at.desc()).all()

@router.delete("/{document_id}")
def delete_document(document_id: UUID, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    document = db.query(Document).filter(Document.id == document_id, Document.user_id == current_user.id).first()
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
        
    # Remove vectors
    delete_document_vectors(document_id)
    
    # Remove file
    if os.path.exists(document.file_path):
        try:
            os.remove(document.file_path)
        except:
            pass
            
    db.delete(document)
    db.commit()
    return {"status": "success", "message": "Document deleted"}

@router.post("/{document_id}/summarize", response_model=DocumentSummaryResponse)
def summarize_document(document_id: UUID, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    document = db.query(Document).filter(Document.id == document_id, Document.user_id == current_user.id).first()
    if not document:
        raise HTTPException(status_code=404, detail="Document not found")
    
    if document.summary:
        return {"summary": document.summary}
        
    summary = generate_document_summary(db, document)
    return {"summary": summary}

@router.post("/{document_id}/quiz")
def generate_quiz_from_document(document_id: UUID, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    # Returns quiz_id
    quiz = generate_quiz(db, current_user.id, "", 10, document_id)
    return {"status": "success", "data": {"quiz_id": quiz.id}}

from fastapi.responses import FileResponse

@router.get("/{document_id}/download")
def download_document(document_id: UUID, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    document = db.query(Document).filter(Document.id == document_id, Document.user_id == current_user.id).first()
    if not document or not os.path.exists(document.file_path):
        raise HTTPException(status_code=404, detail="Document not found")
    return FileResponse(path=document.file_path, filename=document.original_name)

@router.get("/{document_id}/content")
def get_document_raw_content(document_id: UUID, db: Session = Depends(get_db), current_user: User = Depends(get_current_user)):
    document = db.query(Document).filter(Document.id == document_id, Document.user_id == current_user.id).first()
    if not document or not os.path.exists(document.file_path):
        raise HTTPException(status_code=404, detail="Document not found")
    
    # Extract text to show to user
    text = ""
    html = ""
    ext = os.path.splitext(document.file_path)[1].lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(document.file_path)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        elif ext == ".docx":
            try:
                import mammoth
                with open(document.file_path, "rb") as docx_file:
                    result = mammoth.convert_to_html(docx_file)
                    html = result.value
            except Exception as e:
                print(f"Mammoth error: {e}")
                from docx import Document as DocxDocument
                doc = DocxDocument(document.file_path)
                for para in doc.paragraphs:
                    if para.text.strip():
                        text += para.text + "\n"
        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(document.file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
        else:
            with open(document.file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
    except Exception:
        text = "Lỗi: Không thể trích xuất văn bản từ tài liệu này để hiển thị."
        
    return {"text": text, "html": html}
